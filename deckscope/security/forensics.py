"""File-level forensics: text a human never sees, but a text extractor does.

Plain-text extraction flattens a deck. Anything hidden by *rendering* — white text
on a white background, 1pt fonts, shapes parked off the slide, layers behind
images, hidden slides, speaker notes — arrives at the model looking exactly like
the headline. These scanners re-open the original file and recover what the
extractor threw away.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from .policy import SecurityPolicy
from .report import Finding, ScanReport
from .text_scanner import scan_text

# ---------------------------------------------------------------- colour utils


def _luminance(rgb: Tuple[float, float, float]) -> float:
    """Relative luminance, 0 (black) to 1 (white). Inputs are 0-1 floats."""
    r, g, b = [max(0.0, min(1.0, float(c))) for c in rgb]
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _to_rgb(value: Any) -> Optional[Tuple[float, float, float]]:
    """Normalize the several colour shapes PDF/PPTX libraries hand back."""
    if value is None:
        return None
    if isinstance(value, (int, float)):
        g = float(value)
        g = g / 255.0 if g > 1 else g
        return (g, g, g)
    if isinstance(value, (list, tuple)):
        vals = [float(v) for v in value if isinstance(v, (int, float))]
        if len(vals) == 1:
            return _to_rgb(vals[0])
        if len(vals) == 3:
            if max(vals) > 1:
                vals = [v / 255.0 for v in vals]
            return (vals[0], vals[1], vals[2])
        if len(vals) == 4:  # CMYK
            c, m, y, k = vals
            return ((1 - c) * (1 - k), (1 - m) * (1 - k), (1 - y) * (1 - k))
    if isinstance(value, str) and len(value) in (6, 8):
        try:
            return (int(value[0:2], 16) / 255, int(value[2:4], 16) / 255,
                    int(value[4:6], 16) / 255)
        except ValueError:
            return None
    return None


# ---------------------------------------------------------------- PDF

def scan_pdf(path: Path, policy: SecurityPolicy) -> ScanReport:
    """Recover invisible-by-rendering text from a PDF and screen it."""
    rep = ScanReport(target=f"deck forensics ({path.name})")
    try:
        import pdfplumber
    except ImportError:
        rep.add(Finding("forensics_unavailable", "low", str(path.name),
                        "pdfplumber is not installed, so hidden-text forensics were "
                        "skipped for this PDF. Install it with: pip install pdfplumber"))
        return rep

    try:
        pdf = pdfplumber.open(str(path))
    except Exception as exc:  # noqa: BLE001
        rep.add(Finding("forensics_error", "low", str(path.name),
                        f"Could not open the PDF for forensics: {exc}"))
        return rep

    with pdf:
        if policy.scan_metadata:
            rep.extend(_scan_metadata(dict(pdf.metadata or {}), f"{path.name} metadata"))

        for pageno, page in enumerate(pdf.pages, 1):
            rep.scanned_items += 1
            try:
                chars = page.chars or []
            except Exception:  # noqa: BLE001
                continue
            if not chars:
                continue

            page_bg = _page_background(page)
            buckets: Dict[str, List[Any]] = {"invisible": [], "tiny": [], "offpage": []}

            for ch in chars:
                size = float(ch.get("size") or 0)
                if size and size < policy.min_font_pt:
                    buckets["tiny"].append(ch)
                    continue
                if ch.get("non_stroking_color") is not None:
                    rgb = _to_rgb(ch.get("non_stroking_color"))
                    if rgb is not None:
                        if abs(_luminance(rgb) - page_bg) < policy.contrast_threshold:
                            buckets["invisible"].append(ch)
                            continue
                # render mode 3 = "invisible" (used legitimately by OCR layers, but
                # also the simplest way to hide a payload)
                if ch.get("render_mode") == 3 or ch.get("upright") is False:
                    buckets["invisible"].append(ch)
                    continue
                x0, top = float(ch.get("x0", 0)), float(ch.get("top", 0))
                if x0 < -5 or top < -5 or x0 > float(page.width) + 5 or \
                        top > float(page.height) + 5:
                    buckets["offpage"].append(ch)

            for kind, label, sev in (
                ("invisible", "text whose colour matches the background", "high"),
                ("tiny", f"text smaller than {policy.min_font_pt}pt", "high"),
                ("offpage", "text positioned outside the page boundary", "high"),
            ):
                items = buckets[kind]
                if len(items) < 12:   # a stray glyph is noise, a sentence is not
                    continue
                hidden = "".join(c.get("text", "") for c in items).strip()
                where = f"page {pageno}"
                sub = scan_text(hidden, where)
                worst = sub.risk
                rep.add(Finding(
                    code=f"hidden_{kind}", severity=sev if worst == "clean" else "critical",
                    where=where,
                    detail=(f"{len(items)} characters of {label} — invisible to a human "
                            f"reader but delivered to the AI."
                            + (f" It contains AI-directed instructions "
                               f"({', '.join(sorted({f.code for f in sub.findings}))})."
                               if sub.findings else
                               " No injection language detected in it, but it is still "
                               "content the reader cannot see.")),
                    excerpt=hidden[:280],
                    action="redacted" if policy.should_redact(
                        "critical" if sub.findings else sev) else "flagged"))
                rep.extend(sub)
    return rep


def _page_background(page: Any) -> float:
    """Best guess at page luminance: a full-page light rect, else assume white."""
    try:
        for rect in (page.rects or []):
            w = float(rect.get("width") or 0)
            h = float(rect.get("height") or 0)
            if w > float(page.width) * 0.9 and h > float(page.height) * 0.9:
                rgb = _to_rgb(rect.get("non_stroking_color"))
                if rgb is not None:
                    return _luminance(rgb)
    except Exception:  # noqa: BLE001
        pass
    return 1.0


# ---------------------------------------------------------------- PPTX

def scan_pptx(path: Path, policy: SecurityPolicy) -> ScanReport:
    """Recover hidden runs, off-slide shapes, hidden slides, and notes from a PPTX."""
    rep = ScanReport(target=f"deck forensics ({path.name})")
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        rep.add(Finding("forensics_unavailable", "low", path.name,
                        "python-pptx is not installed, so hidden-text forensics were "
                        "skipped. Install it with: pip install python-pptx"))
        return rep

    try:
        prs = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        rep.add(Finding("forensics_error", "low", path.name,
                        f"Could not open the deck for forensics: {exc}"))
        return rep

    if policy.scan_metadata:
        cp = prs.core_properties
        meta = {k: getattr(cp, k, None) for k in
                ("title", "subject", "keywords", "comments", "category", "author")}
        rep.extend(_scan_metadata(meta, f"{path.name} metadata"))

    sw, sh = int(prs.slide_width or 0), int(prs.slide_height or 0)

    for i, slide in enumerate(prs.slides, 1):
        rep.scanned_items += 1
        where = f"slide {i}"

        if _slide_hidden(slide):
            text = _shape_text(slide.shapes)
            if text.strip():
                sub = scan_text(text, where)
                rep.add(Finding(
                    "hidden_slide", "critical" if sub.findings else "medium", where,
                    "This slide is marked hidden — it never appears when the deck is "
                    "presented, but its text still reaches the AI."
                    + (" It contains AI-directed instructions." if sub.findings else ""),
                    excerpt=text[:280],
                    action="redacted" if policy.should_redact("critical") else "flagged"))
                rep.extend(sub)

        bg_lum = _slide_bg_luminance(slide)

        for shape in slide.shapes:
            # off-slide placement
            try:
                if sw and sh and shape.left is not None and shape.top is not None:
                    left, top = int(shape.left), int(shape.top)
                    if (left + int(shape.width or 0) < 0 or top + int(shape.height or 0) < 0
                            or left > sw or top > sh):
                        text = _one_shape_text(shape)
                        if text.strip():
                            sub = scan_text(text, where)
                            rep.add(Finding(
                                "offslide_shape",
                                "critical" if sub.findings else "high", where,
                                "A text box is positioned outside the visible slide area "
                                "— invisible when presented, but extracted as text."
                                + (" It contains AI-directed instructions."
                                   if sub.findings else ""),
                                excerpt=text[:280],
                                action="redacted" if policy.should_redact("high") else "flagged"))
                            rep.extend(sub)
                            continue
            except Exception:  # noqa: BLE001
                pass

            if not getattr(shape, "has_text_frame", False):
                continue
            shape_lum = _shape_fill_luminance(shape)
            base_lum = shape_lum if shape_lum is not None else bg_lum

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    body = (run.text or "").strip()
                    if len(body) < 15:
                        continue
                    reasons = []
                    try:
                        size_pt = run.font.size.pt if run.font.size else None
                    except Exception:  # noqa: BLE001
                        size_pt = None
                    if size_pt is not None and size_pt < policy.min_font_pt:
                        reasons.append(f"font size {size_pt}pt")
                    rgb = _run_rgb(run)
                    if rgb is not None and abs(_luminance(rgb) - base_lum) < policy.contrast_threshold:
                        reasons.append("text colour matches the background")
                    try:
                        if run.font.color and getattr(run.font.color, "_xFill", None) is not None:
                            alpha = _run_alpha(run)
                            if alpha is not None and alpha < 0.15:
                                reasons.append(f"text {int(alpha*100)}% opaque")
                    except Exception:  # noqa: BLE001
                        pass
                    if not reasons:
                        continue
                    sub = scan_text(body, where)
                    rep.add(Finding(
                        "invisible_render",
                        "critical" if sub.findings else "high", where,
                        f"Text a human reader cannot see ({'; '.join(reasons)}), but the "
                        f"AI receives in full."
                        + (" It contains AI-directed instructions." if sub.findings else ""),
                        excerpt=body[:280],
                        action="redacted" if policy.should_redact("high") else "flagged"))
                    rep.extend(sub)

        if policy.scan_speaker_notes:
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text or ""
                    if notes.strip():
                        sub = scan_text(notes, f"{where} speaker notes")
                        if sub.findings:
                            rep.add(Finding(
                                "notes_injection", "critical", f"{where} speaker notes",
                                "Speaker notes contain AI-directed instructions. Notes are "
                                "never shown to an audience, which makes them a favoured "
                                "hiding place.",
                                excerpt=notes[:280],
                                action="redacted" if policy.should_redact("critical") else "flagged"))
                            rep.extend(sub)
            except Exception:  # noqa: BLE001
                pass
    return rep


def _slide_hidden(slide: Any) -> bool:
    try:
        return slide._element.get("show") == "0"
    except Exception:  # noqa: BLE001
        return False


def _shape_text(shapes: Any) -> str:
    parts = []
    for sh in shapes:
        parts.append(_one_shape_text(sh))
    return "\n".join(p for p in parts if p)


def _one_shape_text(shape: Any) -> str:
    try:
        if getattr(shape, "has_text_frame", False):
            return shape.text_frame.text or ""
    except Exception:  # noqa: BLE001
        pass
    return ""


def _run_rgb(run: Any) -> Optional[Tuple[float, float, float]]:
    try:
        color = run.font.color
        if color is None or color.type is None:
            return None
        rgb = getattr(color, "rgb", None)
        if rgb is None:
            return None
        s = str(rgb)
        return (int(s[0:2], 16) / 255, int(s[2:4], 16) / 255, int(s[4:6], 16) / 255)
    except Exception:  # noqa: BLE001
        return None


def _run_alpha(run: Any) -> Optional[float]:
    try:
        xml = run.font.color._xFill.xml
        import re
        m = re.search(r'<a:alpha val="(\d+)"', xml)
        if m:
            return int(m.group(1)) / 100000.0
    except Exception:  # noqa: BLE001
        return None
    return None


def _slide_bg_luminance(slide: Any) -> float:
    try:
        xml = slide._element.xml
        import re
        m = re.search(r'<a:bgFill>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"', xml, re.S)
        if m:
            rgb = _to_rgb(m.group(1))
            if rgb:
                return _luminance(rgb)
    except Exception:  # noqa: BLE001
        pass
    return 1.0


def _shape_fill_luminance(shape: Any) -> Optional[float]:
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # solid
            rgb = fill.fore_color.rgb
            s = str(rgb)
            return _luminance((int(s[0:2], 16) / 255, int(s[2:4], 16) / 255,
                               int(s[4:6], 16) / 255))
    except Exception:  # noqa: BLE001
        return None
    return None


# ---------------------------------------------------------------- DOCX + metadata

def scan_docx(path: Path, policy: SecurityPolicy) -> ScanReport:
    rep = ScanReport(target=f"deck forensics ({path.name})")
    try:
        import docx
    except ImportError:
        return rep
    try:
        d = docx.Document(str(path))
    except Exception as exc:  # noqa: BLE001
        rep.add(Finding("forensics_error", "low", path.name, str(exc)))
        return rep

    if policy.scan_metadata:
        cp = d.core_properties
        rep.extend(_scan_metadata(
            {k: getattr(cp, k, None) for k in
             ("title", "subject", "keywords", "comments", "category", "author")},
            f"{path.name} metadata"))

    for pi, para in enumerate(d.paragraphs, 1):
        for run in para.runs:
            body = (run.text or "").strip()
            if len(body) < 15:
                continue
            reasons = []
            try:
                if run.font.size and run.font.size.pt < policy.min_font_pt:
                    reasons.append(f"font size {run.font.size.pt}pt")
            except Exception:  # noqa: BLE001
                pass
            try:
                if run.font.hidden:
                    reasons.append("marked hidden")
            except Exception:  # noqa: BLE001
                pass
            try:
                rgb = run.font.color.rgb if run.font.color else None
                if rgb is not None and str(rgb).upper() == "FFFFFF":
                    reasons.append("white text")
            except Exception:  # noqa: BLE001
                pass
            if reasons:
                sub = scan_text(body, f"paragraph {pi}")
                rep.add(Finding(
                    "invisible_render", "critical" if sub.findings else "high",
                    f"paragraph {pi}",
                    f"Text a reader cannot see ({'; '.join(reasons)}).",
                    excerpt=body[:280],
                    action="redacted" if policy.should_redact("high") else "flagged"))
                rep.extend(sub)
    return rep


def _scan_metadata(meta: Dict[str, Any], where: str) -> ScanReport:
    """Document metadata is invisible in every viewer and extracted by many parsers."""
    rep = ScanReport(target=where)
    for key, value in (meta or {}).items():
        if not value or not isinstance(value, str) or len(value) < 15:
            continue
        sub = scan_text(value, f"{where}: {key}")
        if sub.findings:
            rep.add(Finding(
                "metadata_injection", "critical", f"{where}: {key}",
                f"The document's `{key}` metadata field contains AI-directed "
                f"instructions. Metadata is invisible in every normal viewer.",
                excerpt=value[:280], action="redacted"))
            rep.extend(sub)
    return rep


SCANNERS = {".pdf": scan_pdf, ".pptx": scan_pptx, ".docx": scan_docx}


def scan_file(path: str, policy: SecurityPolicy) -> ScanReport:
    p = Path(path)
    fn = SCANNERS.get(p.suffix.lower())
    if not fn or not p.exists():
        return ScanReport(target=f"deck forensics ({p.name})")
    return fn(p, policy)
