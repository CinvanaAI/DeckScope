"""`deckscope improve deck.pdf` — the analysis run backwards: instead of
telling an investor what's wrong with the deck, build the founder the
strongest version of the deck that survives the same audit.

The flow: run the normal analysis (or reuse a finished run), then a single
revision pass rewrites the deck slide by slide against the audit — keep
what held up, correct what the evidence contradicted, add slides for the
blind spots diligence will find, and turn unanswerable questions into
explicit slots only the founder can fill.

The honesty rules are the product, and they are enforced in code, not
requested in the prompt:

- Market and industry corrections may cite ONLY the run's own bibliography
  ([S#]); the citation audit strips anything else and says so.
- A revised or new line asserting a figure with no surviving source is
  demoted to a founder-input slot — DeckScope does not invent numbers in
  either direction, and a deck-builder that fabricates traction is a
  worse product than no deck-builder.
- A kept line that token-matches a contested claim is flagged "kept
  against the evidence" — the founder may keep it, but the blueprint
  says an investor-side run will catch it.

Because the loader ingests .txt and .md, `deckscope improve notes.txt`
is the build-from-scratch path: raw founder notes in, a structured,
evidence-checked deck blueprint out — same command, same rules.
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from .audit_report import _FIGURE, _YEAR_ONLY, build_registry
from .diff import _overlap, _tokens

_KINDS = ("kept", "revised", "new", "founder-input")
#: Above this token overlap a kept line and a contested claim are the same
#: assertion. Deliberately below diff's pairing bar: a founder lightly
#: rewording a contested claim should still be warned.
KEPT_MATCH = 0.45


# ------------------------------------------------------------ the brief

def build_brief(comparison: Dict[str, Any], registry: Any) -> str:
    """Everything the reviser is allowed to know, assembled by code so the
    prompt cannot drift from what the run established."""
    L: List[str] = []
    add = L.append
    audit = [r for r in (comparison.get("claim_audit") or [])
             if isinstance(r, dict)]

    contested = [r for r in audit if str(r.get("assessment", "")).lower()
                 in ("contradicted", "partially-supported")]
    if contested:
        add("CLAIMS THE EVIDENCE CONTRADICTS — correct these to what the "
            "cited sources actually show:")
        for r in contested:
            add(f"- ({r.get('id')}) {r.get('claim', '')}")
            if r.get("market_evidence"):
                sids = " ".join(f"[{s}]" for s in (r.get("source_ids") or []))
                add(f"  evidence: {r['market_evidence']} {sids}".rstrip())
            if r.get("delta"):
                add(f"  gap: {r['delta']}")
        add("")

    supported = [r for r in audit
                 if str(r.get("assessment", "")).lower() == "supported"]
    if supported:
        add("CLAIMS THAT SURVIVED THE AUDIT — keep these, they are "
            "strengths:")
        for r in supported:
            add(f"- ({r.get('id')}) {r.get('claim', '')}")
        add("")

    unverifiable = [r for r in audit
                    if str(r.get("assessment", "")).lower() == "unverifiable"]
    if unverifiable:
        add("CLAIMS ONLY THE FOUNDER CAN SETTLE — turn each into a "
            "founder-input line with a bracketed slot; do not assert them "
            "as verified:")
        for r in unverifiable:
            add(f"- ({r.get('id')}) {r.get('claim', '')}")
        add("")

    blind = [b for b in ((comparison.get("alignment") or {})
                         .get("blind_spots") or []) if isinstance(b, dict)]
    if blind:
        add("WHAT THE DECK NEVER ADDRESSES — investor-side diligence found "
            "these; add a slide or line that addresses each head-on:")
        for b in blind:
            add(f"- {b.get('what', '')} — {b.get('why_it_matters', '')}")
        add("")

    questions = [str(q).strip() for q in (comparison.get("questions") or [])
                 if str(q).strip()]
    if questions:
        add("QUESTIONS AN INVESTOR WILL ASK — answer in the deck where the "
            "founder can, so the deck answers before the room asks:")
        for q in questions:
            add(f"- {q}")
        add("")

    add("BIBLIOGRAPHY — the ONLY citable sources:")
    add(registry.prompt_block(char_budget=45_000))
    return "\n".join(L)


# ---------------------------------------------------- validation (the law)

def _has_figure(text: str) -> bool:
    return any(not _YEAR_ONLY.fullmatch(m.strip())
               for m in _FIGURE.findall(text or ""))


def validate_revision(rev: Dict[str, Any], registry: Any,
                      comparison: Dict[str, Any]) -> Tuple[Dict[str, Any],
                                                           Dict[str, Any]]:
    """Deterministic enforcement after the model call. Returns the cleaned
    revision plus notes on everything that was enforced."""
    from ..sources import audit_fragment  # local: package-relative depth

    notes = {"stripped_citations": 0, "demoted_lines": 0,
             "kept_against_evidence": 0}

    fa = audit_fragment(rev, registry, strip=True, where="revision")
    notes["stripped_citations"] = (len(fa.dangling) + len(fa.quarantined)
                                   + len(fa.unadmitted))

    contested = [(str(r.get("id", "")), _tokens(r.get("claim", "")))
                 for r in (comparison.get("claim_audit") or [])
                 if isinstance(r, dict)
                 and str(r.get("assessment", "")).lower()
                 in ("contradicted", "partially-supported")]

    slides = []
    for slide in (rev.get("slides") or []):
        if not isinstance(slide, dict):
            continue
        lines = []
        for line in (slide.get("lines") or []):
            if not isinstance(line, dict) or not str(
                    line.get("text", "")).strip():
                continue
            kind = str(line.get("kind", "")).strip().lower()
            if kind not in _KINDS:
                kind = "revised"
            sids = [s for s in (line.get("source_ids") or []) if s]

            # The core rule: a new/revised figure with no surviving source
            # becomes the founder's problem to source, visibly — never a
            # fact the blueprint invented.
            if kind in ("revised", "new") and _has_figure(
                    line.get("text", "")) and not sids:
                kind = "founder-input"
                line["demoted"] = True
                notes["demoted_lines"] += 1
                fi = rev.setdefault("founder_inputs", [])
                fi.append({"slot": line["text"],
                           "what_to_provide": "a source for this figure, "
                           "or your real number",
                           "standard": "a number with a named source or "
                           "your own books behind it"})

            if kind == "founder-input" and "[" not in line.get("text", ""):
                line["text"] = (str(line["text"]).rstrip(". ")
                                + " — [YOUR NUMBER/DETAIL HERE]")

            if kind == "kept":
                toks = _tokens(line.get("text", ""))
                for cid, ctoks in contested:
                    if _overlap(toks, ctoks) >= KEPT_MATCH:
                        line["kept_against_evidence"] = cid
                        notes["kept_against_evidence"] += 1
                        break

            line["kind"] = kind
            line["source_ids"] = sids
            lines.append(line)
        if lines:
            slide["lines"] = lines
            slides.append(slide)

    for i, slide in enumerate(slides, 1):
        slide["n"] = i
    rev["slides"] = slides
    return rev, notes


# ------------------------------------------------------------- rendering

_MARK = {"revised": "· revised", "new": "· new", "founder-input": "✎ you"}


def render_blueprint(rev: Dict[str, Any], notes: Dict[str, Any],
                     registry: Any, source_label: str) -> str:
    L: List[str] = []
    add = L.append
    company = rev.get("company") or "your company"
    add(f"# Improved deck — {company}")
    add("")
    add(f"*Rebuilt from the audit of {source_label}. Every market "
        "correction cites the run's own retrieved sources ([S#] — full "
        "list in the appendix). Lines marked \"✎ you\" are slots only you "
        "can fill: DeckScope will not invent your numbers, in either "
        "direction. A kept line flagged ⚠ is one an investor-side run of "
        "this same tool will contest — keep it knowingly or fix it.*")
    add("")
    if rev.get("positioning"):
        add(f"**The story, in one line:** {rev['positioning']}")
        add("")
    enforced = []
    if notes.get("stripped_citations"):
        enforced.append(f"{notes['stripped_citations']} citation(s) the "
                        "model offered were not in the run's bibliography "
                        "and were removed")
    if notes.get("demoted_lines"):
        enforced.append(f"{notes['demoted_lines']} figure line(s) arrived "
                        "with no source and were demoted to founder slots")
    if enforced:
        add(f"*Enforced during validation: {'; '.join(enforced)}.*")
        add("")

    for slide in rev.get("slides") or []:
        add(f"## Slide {slide.get('n')} — {slide.get('title', '')}")
        add("")
        if slide.get("purpose"):
            add(f"*{slide['purpose']}*")
            add("")
        for line in slide.get("lines") or []:
            kind = line.get("kind")
            sids = " ".join(f"[{s}]" for s in (line.get("source_ids") or []))
            mark = _MARK.get(kind, "")
            row = f"- {line.get('text', '')}"
            if sids:
                row += f" {sids}"
            if mark:
                row += f"  *({mark}"
                if line.get("because"):
                    row += f" — {line['because']}"
                row += ")*"
            if line.get("kept_against_evidence"):
                row += (f"  ⚠ *kept against the evidence "
                        f"({line['kept_against_evidence']}) — an "
                        f"investor-side run will contest this line*")
            add(row)
        if slide.get("speaker_note"):
            add("")
            add(f"> Speaker note: {slide['speaker_note']}")
        add("")

    cut = [c for c in (rev.get("cut") or []) if isinstance(c, dict)]
    if cut:
        add("## Cut from the deck")
        add("")
        for c in cut:
            add(f"- **{c.get('what', '')}** — {c.get('why', '')}")
        add("")

    fi = [f for f in (rev.get("founder_inputs") or []) if isinstance(f, dict)]
    if fi:
        add("## Only you can fill these")
        add("")
        seen = set()
        for f in fi:
            key = str(f.get("slot", ""))[:80]
            if key in seen:
                continue
            seen.add(key)
            add(f"- **{f.get('slot', '')}** — {f.get('what_to_provide', '')}")
            if f.get("standard"):
                add(f"  - A good answer: {f['standard']}")
        add("")

    cited = []
    seen_s = set()
    for slide in rev.get("slides") or []:
        for line in slide.get("lines") or []:
            for s in line.get("source_ids") or []:
                if s not in seen_s:
                    seen_s.add(s)
                    cited.append(s)
    if cited:
        add("## Evidence appendix — carry these into your data room")
        add("")
        for s in cited:
            src = registry.find(s)
            if src is not None:
                add(f"- **[{s}]** {src.title or '(untitled)'} — "
                    f"{src.url or 'n/a'}")
        add("")

    add("---")
    add("*An AI-assisted draft built from an evidence audit — not a "
        "finished deck and not investment or fundraising advice. Verify "
        "every figure before it goes in front of an investor; the "
        "bracketed slots are yours to fill, and the flagged lines are "
        "yours to defend.*")
    return "\n".join(L)


def write_pptx(rev: Dict[str, Any], path: Path) -> Optional[Path]:
    """A plain, editable starting deck — not a designed one. Returns None
    (with no file) when python-pptx is unavailable."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return None

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = rev.get("company") or "Deck"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = rev.get("positioning") or ""

    for slide in rev.get("slides") or []:
        ps = prs.slides.add_slide(body_layout)
        ps.shapes.title.text = str(slide.get("title", ""))[:90]
        body = ps.placeholders[1].text_frame
        first = True
        for line in slide.get("lines") or []:
            text = str(line.get("text", ""))
            if line.get("kind") == "founder-input":
                text = "▢ " + text
            sids = " ".join(f"[{s}]" for s in (line.get("source_ids") or []))
            if sids:
                text += f" {sids}"
            para = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            para.text = text
            para.font.size = Pt(18)
        because = [f"{ln.get('kind')}: {ln.get('because')}"
                   for ln in (slide.get("lines") or []) if ln.get("because")]
        note = "\n".join(([slide.get("speaker_note")] if
                          slide.get("speaker_note") else []) + because)
        if note:
            ps.notes_slide.notes_text_frame.text = note

    srcs = prs.slides.add_slide(body_layout)
    srcs.shapes.title.text = "Sources"
    tf = srcs.placeholders[1].text_frame
    first = True
    seen = set()
    for slide in rev.get("slides") or []:
        for line in slide.get("lines") or []:
            for sid in line.get("source_ids") or []:
                if sid in seen:
                    continue
                seen.add(sid)
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                para.text = f"[{sid}]"
                para.font.size = Pt(14)
    prs.save(path)
    return path


# ---------------------------------------------------------------- command

def _reviser(provider, cache_dir, verbose):
    from ..agents.base import Agent
    from ..prompts.templates import DECK_REVISE_SYSTEM, DECK_REVISE_USER
    from ..schemas import REVISION_SCHEMA, coerce, schema_block

    class DeckReviser(Agent):
        name = "revise"
        label = "Deck Reviser"

        def run(self, brief: str, deck_text: str) -> Dict[str, Any]:
            user = DECK_REVISE_USER.format(
                schema=schema_block(REVISION_SCHEMA, "DeckRevision"),
                brief=brief, deck_text=deck_text[:60_000])
            self.emit("rebuilding the deck against the audit")
            result = self.cached_json(
                self.cache_key(brief=brief, deck=deck_text[:60_000]),
                lambda: self.complete_json(DECK_REVISE_SYSTEM, user))
            return coerce(result, REVISION_SCHEMA)

    return DeckReviser(provider, cache_dir=cache_dir, verbose=verbose)


def command(args: Any) -> int:
    import sys

    from ..console import out as _out

    def _err(msg):
        _out(msg, file=sys.stderr)

    from ..config import (OutputConfig, ProviderConfig, ResearchConfig,
                           RunConfig, load_config)
    from ..ingest.loader import load_deck
    from ..orchestrator import Pipeline
    from ..providers.registry import get_provider
    from ..security.screening import screen_deck
    from ..tiering import NDAGuard, is_local

    demo = bool(getattr(args, "demo", False))
    nda = bool(getattr(args, "nda", False))
    lens = getattr(args, "lens", None) or "founder"

    if demo:
        here = Path(__file__).resolve().parent.parent
        deck_path = here / "examples" / "sample_deck.md"
        if not deck_path.exists():
            _err(f"The packaged sample deck is missing ({deck_path}) — "
                 "reinstall DeckScope, or run improve on your own file.")
            return 2
        cfg = RunConfig(deck_path=str(deck_path), lenses=[lens],
                        provider=ProviderConfig(name="mock"),
                        research=ResearchConfig(name="none"),
                        output=OutputConfig(formats=["markdown"]),
                        cache_dir=None)
    else:
        if not getattr(args, "deck", None):
            _err("Give me a deck file (or a .txt/.md of raw notes — "
                 "improve builds a deck from those too).")
            return 2
        deck_path = Path(args.deck)
        if not deck_path.exists():
            _err(f"Not found: {deck_path}")
            return 2
        cfg = load_config(getattr(args, "config", None))
        if getattr(args, "provider", None):
            cfg.provider.name = args.provider
        cfg.deck_path = str(deck_path)
        cfg.lenses = [lens]
        if nda:
            # Deck-derived search queries must not leave the machine.
            cfg.research.name = "none"
        cfg.__post_init__()

    if nda:
        # BOTH providers: the pipeline sends the full deck to
        # `extract_provider` when one is configured (sixth external audit —
        # a local main model with a hosted extraction model passed the old
        # single-provider check and shipped the deck out anyway).
        for label, pc in (("model", cfg.provider),
                          ("extraction model", cfg.extract_provider)):
            if pc is not None and not is_local(pc):
                _err(f"--nda refused: the configured {label} "
                     f"('{pc.name}') is not local, and improve sends the "
                     "deck and the audit to it. Use a local model or drop "
                     "--nda.")
                return 4

    out_dir = Path(getattr(args, "out", None) or "deckscope_output")
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = deck_path.stem

    # ---- the forward run (or a finished one, reused)
    from_run = getattr(args, "from_run", None)
    if from_run:
        try:
            rec = json.loads(Path(from_run).read_text(encoding="utf-8"))
        except (OSError, ValueError) as e:
            _err(f"Could not read the saved run: {e}")
            return 2
        comparisons = rec.get("comparisons") or {}
        if lens not in comparisons and comparisons:
            lens = next(iter(comparisons))
        comparison = comparisons.get(lens) or {}
        try:
            registry = build_registry(
                [s for s in (rec.get("sources") or []) if isinstance(s, dict)])
        except ValueError as e:
            _err(f"The saved run's sources did not load: {e}")
            return 2
        provider = get_provider(cfg.provider)
        source_label = f"{deck_path.name} (audit reused from {Path(from_run).name})"
    else:
        pipe = Pipeline(cfg)
        provider = pipe.provider
        corpus = None
        if demo:
            from ..cli import _demo_corpus
            corpus = _demo_corpus(Path(__file__).resolve().parent.parent)
        try:
            result = pipe.run(corpus=corpus)
        except Exception as e:  # noqa: BLE001 — surfaced, not swallowed
            _err(f"The analysis run failed before any revision happened: "
                 f"{type(e).__name__}: {e}")
            return 1
        comparison = result.comparisons.get(lens) or {}
        registry = result.registry
        source_label = deck_path.name

    if not (comparison.get("claim_audit") or []):
        _err("The run produced no claim audit to rebuild from — nothing "
             "honest to improve against. Run with research enabled.")
        return 1

    # ---- the deck text the reviser sees (screened like any other input)
    doc = load_deck(str(deck_path))
    doc, _scan = screen_deck(doc, cfg.security, deck_path=str(deck_path))
    guard = NDAGuard(enabled=nda)
    guard.protect(doc.text)

    brief = build_brief(comparison, registry)
    reviser = _reviser(provider, cfg.cache_dir, cfg.verbose)
    rev = reviser.run(brief, doc.text)
    rev, notes = validate_revision(rev, registry, comparison)

    if not rev.get("slides"):
        _err("The reviser returned no usable slides — nothing was written.")
        return 1

    md_path = out_dir / f"{slug}_{lens}_improved.md"
    md_path.write_text(
        render_blueprint(rev, notes, registry, source_label),
        encoding="utf-8")

    files = [str(md_path)]
    if getattr(args, "pptx", False):
        p = write_pptx(rev, out_dir / f"{slug}_improved.pptx")
        if p is None:
            _out("(.pptx skipped — python-pptx is not installed; "
                 "the markdown blueprint has everything.)")
        else:
            files.append(str(p))

    n_lines = sum(len(s.get("lines") or []) for s in rev["slides"])
    kinds = {k: 0 for k in _KINDS}
    for s in rev["slides"]:
        for ln in s.get("lines") or []:
            kinds[ln.get("kind", "revised")] = kinds.get(
                ln.get("kind", "revised"), 0) + 1
    _out(f"\n{len(rev['slides'])} slide(s), {n_lines} line(s): "
         f"{kinds['kept']} kept, {kinds['revised']} revised, "
         f"{kinds['new']} new, {kinds['founder-input']} for you to fill.")
    if notes["stripped_citations"] or notes["demoted_lines"]:
        _out(f"Enforced: {notes['stripped_citations']} invalid citation(s) "
             f"stripped, {notes['demoted_lines']} unsourced figure(s) "
             f"demoted to founder slots.")
    if notes["kept_against_evidence"]:
        _out(f"⚠ {notes['kept_against_evidence']} kept line(s) contradict "
             "the evidence — flagged in the blueprint.")
    for f in files:
        _out(f"Written: {f}")
    return 0
