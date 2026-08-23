"""Slide output — the summary deck you present back to a committee or a founder."""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List

from .common import ASSESSMENT_WORD, as_list, header_block, score_color, theme as get_theme, txt


def render(result, out_dir: Path, base: str, theme: str = "slate", **kw: Any) -> List[str]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        raise RuntimeError("Slide output needs python-pptx: pip install python-pptx") from None

    t = get_theme(theme)

    def rgb(hexstr: str) -> Any:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    paths = []
    for lens, comp in result.comparisons.items():
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        W, H = prs.slide_width, prs.slide_height
        blank = prs.slide_layouts[6]
        h = header_block(result, lens)

        def new_slide(title: str = "", subtitle: str = ""):
            s = prs.slides.add_slide(blank)
            bg = s.background.fill
            bg.solid()
            bg.fore_color.rgb = rgb(t["bg"])
            if title:
                box = s.shapes.add_textbox(Inches(0.7), Inches(0.45),
                                           W - Inches(1.4), Inches(0.9))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = title
                r.font.size, r.font.bold, r.font.color.rgb = Pt(30), True, rgb(t["ink"])
                if subtitle:
                    p2 = tf.add_paragraph()
                    r2 = p2.add_run()
                    r2.text = subtitle
                    r2.font.size, r2.font.color.rgb = Pt(13), rgb(t["muted"])
                line = s.shapes.add_shape(1, Inches(0.7), Inches(1.55),
                                          W - Inches(1.4), Pt(2.5))
                line.fill.solid()
                line.fill.fore_color.rgb = rgb(t["accent"])
                line.line.fill.background()
            return s

        def body_box(s, top=Inches(1.9), height=None):
            box = s.shapes.add_textbox(Inches(0.7), top, W - Inches(1.4),
                                       height or (H - top - Inches(0.6)))
            box.text_frame.word_wrap = True
            return box.text_frame

        def bullet(tf, text: str, size=14, bold=False, color=None, space=8, level=0):
            p = tf.paragraphs[0] if (not tf.text and len(tf.paragraphs) == 1
                                     and not tf.paragraphs[0].runs) else tf.add_paragraph()
            p.level = level
            p.space_after = Pt(space)
            r = p.add_run()
            r.text = text
            r.font.size, r.font.bold = Pt(size), bold
            r.font.color.rgb = rgb(color or t["ink"])
            return p

        # ---------------------------------------------------------- title
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = rgb(t["accent"])
        box = s.shapes.add_textbox(Inches(1.0), Inches(2.1), W - Inches(2.0), Inches(3.4))
        tf = box.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = result.company
        r.font.size, r.font.bold, r.font.color.rgb = Pt(48), True, rgb("FFFFFF")
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = "Pitch deck claims measured against market evidence"
        r.font.size, r.font.color.rgb = Pt(18), rgb("FFFFFF")
        p = tf.add_paragraph()
        p.space_before = Pt(22)
        r = p.add_run()
        r.text = f"{h['lens']}  ·  {h['verdict']}  ·  {h['score']}/100"
        r.font.size, r.font.bold, r.font.color.rgb = Pt(16), True, rgb("FFFFFF")

        # -------------------------------------------------------- headline
        if h["headline"]:
            s = new_slide("The one-line version")
            tf = body_box(s, Inches(2.4))
            bullet(tf, h["headline"], size=26, bold=True, color=t["accent"], space=18)
            bullet(tf, f"Confidence: {h['confidence']} — "
                       f"{(comp.get('verdict') or {}).get('confidence_rationale', '')}",
                   size=13, color=t["muted"])

        # ------------------------------------------------------- scorecard
        rows = comp.get("scorecard") or []
        if rows:
            s = new_slide("Scorecard", f"Weighted total: {h['score']}/100")
            top = Inches(2.0)
            row_h = Inches(0.52)
            for i, r_ in enumerate(rows[:8]):
                y = top + row_h * i
                lbl = s.shapes.add_textbox(Inches(0.7), y, Inches(3.4), row_h)
                ltf = lbl.text_frame
                ltf.word_wrap = True
                run = ltf.paragraphs[0].add_run()
                run.text = txt(r_.get("dimension"))
                run.font.size, run.font.bold, run.font.color.rgb = Pt(12.5), True, rgb(t["ink"])

                try:
                    pct = max(0.0, min(1.0, float(r_.get("score") or 0) / 10))
                except (TypeError, ValueError):
                    pct = 0.0
                track_w = Inches(5.0)
                track = s.shapes.add_shape(5, Inches(4.2), y + Inches(0.12), track_w, Inches(0.22))
                track.fill.solid()
                track.fill.fore_color.rgb = rgb(t["line"])
                track.line.fill.background()
                if pct > 0:
                    barw = Emu(int(track_w * pct))
                    bar = s.shapes.add_shape(5, Inches(4.2), y + Inches(0.12), barw, Inches(0.22))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = rgb(score_color(r_.get("score"), t))
                    bar.line.fill.background()

                sc = s.shapes.add_textbox(Inches(9.4), y, Inches(1.1), row_h)
                run = sc.text_frame.paragraphs[0].add_run()
                run.text = f"{txt(r_.get('score'))}/10"
                run.font.size, run.font.bold, run.font.color.rgb = Pt(12.5), True, rgb(t["ink"])

                wt = s.shapes.add_textbox(Inches(10.5), y, Inches(2.1), row_h)
                run = wt.text_frame.paragraphs[0].add_run()
                run.text = f"weight {txt(r_.get('weight'))}"
                run.font.size, run.font.color.rgb = Pt(11), rgb(t["muted"])

        # ----------------------------------------------------- claim audit
        for c in (comp.get("claim_audit") or [])[:6]:
            a = (c.get("assessment") or "").lower()
            s = new_slide(f"{txt(c.get('id'))} · {txt(c.get('claim'))}",
                          ASSESSMENT_WORD.get(a, a))
            tf = body_box(s)
            bullet(tf, "What the market evidence shows", size=12, bold=True,
                   color=t["muted"], space=4)
            bullet(tf, txt(c.get("market_evidence")), size=14, space=14)
            if c.get("delta"):
                bullet(tf, "The gap", size=12, bold=True, color=t["muted"], space=4)
                bullet(tf, str(c["delta"]), size=14,
                       color=score_color(3 if a == "contradicted" else 6, t), space=14)
            if c.get("so_what"):
                bullet(tf, "So what", size=12, bold=True, color=t["muted"], space=4)
                bullet(tf, str(c["so_what"]), size=14, space=14)
            cites = ", ".join(str(x) for x in as_list(c.get("source_ids"))) or "none cited"
            bullet(tf, f"Sources: {cites}", size=10.5, color=t["muted"])

        # ------------------------------------------------------- alignment
        align = comp.get("alignment") or {}
        if any(align.values()):
            s = new_slide("Where the deck and the market diverge")
            tf = body_box(s)
            for key, title, color in (
                ("where_deck_overstates", "Deck overstates", t["bad"]),
                ("blind_spots", "Blind spots the deck never addresses", t["bad"]),
                ("where_deck_matches_market", "Deck matches the market", t["good"]),
                ("where_deck_understates", "Deck understates its own case", t["warn"]),
            ):
                items = as_list(align.get(key))
                if not items:
                    continue
                bullet(tf, title, size=13, bold=True, color=color, space=4)
                for i in items[:4]:
                    bullet(tf, f"·  {i}", size=12.5, space=3, level=1)
                bullet(tf, "", size=6, space=6)

        # ----------------------------------------------------------- risks
        risks = comp.get("risks") or []
        if risks:
            s = new_slide("Risks")
            tf = body_box(s)
            for r_ in risks[:7]:
                sev = (r_.get("severity") or "").lower()
                color = {"high": t["bad"], "medium": t["warn"]}.get(sev, t["good"])
                bullet(tf, f"{txt(r_.get('risk'))}", size=14, bold=True, color=color, space=2)
                bullet(tf, f"{sev} severity / {txt(r_.get('likelihood'))} likelihood — "
                           f"{txt(r_.get('mitigation_or_test'))}",
                       size=11.5, color=t["muted"], space=12, level=1)

        # ------------------------------------------------------- questions
        qs = as_list(comp.get("questions"))
        acts = comp.get("actions") or []
        if qs or acts:
            s = new_slide("What to do next")
            tf = body_box(s)
            if qs:
                bullet(tf, "Questions this raises", size=13, bold=True,
                       color=t["muted"], space=6)
                for q in qs[:5]:
                    bullet(tf, f"·  {q}", size=13, space=4, level=1)
                bullet(tf, "", size=6, space=10)
            if acts:
                bullet(tf, "Recommended actions", size=13, bold=True,
                       color=t["muted"], space=6)
                for a in acts[:5]:
                    bullet(tf, f"[{txt(a.get('priority'))}]  {txt(a.get('action'))}",
                           size=13, space=4, level=1)

        # ------------------------------------------------------ references
        reg = getattr(result, "registry", None)
        if reg and reg.sources:
            st = reg.stats()
            chunks = [reg.sources[i:i + 12] for i in range(0, len(reg.sources), 12)]
            for idx, chunk in enumerate(chunks, 1):
                sub = (f"{st['total']} retrieved · {st['cited']} cited · "
                       f"{st['quarantined']} dropped by the security screen"
                       if idx == 1 else f"continued ({idx}/{len(chunks)})")
                s = new_slide("References", sub)
                tf = body_box(s)
                for src in chunk:
                    mark = {"cited": "●", "consulted": "○", "quarantined": "✕"}.get(
                        src.status, "○")
                    color = {"cited": t["ink"], "quarantined": t["bad"]}.get(
                        src.status, t["muted"])
                    bullet(tf, f"{mark} [{src.sid}] {(src.title or src.url or '')[:95]}",
                           size=11, color=color, space=1)
                    if src.url:
                        bullet(tf, f"     {src.url[:110]}", size=9, color=t["muted"],
                               space=5)
        else:
            s = new_slide("References", "No external sources were retrieved")
            tf = body_box(s)
            bullet(tf, "This analysis ran without a web-research backend. Every statement "
                       "rests on the model's training knowledge and on the deck itself, "
                       "and should be treated as unverified.", size=14, color=t["warn"])

        # -------------------------------------------------------- security
        sec = result.security or {}
        if sec:
            risk = sec.get("overall_risk", "clean")
            s = new_slide("Input integrity screen",
                          f"Overall risk: {risk.upper()} · mode {sec.get('mode', 'balanced')}")
            tf = body_box(s)
            if risk == "clean":
                bullet(tf, "The pitch deck and every web source were screened for content "
                           "written to influence the AI rather than inform a human reader "
                           "— hidden text, invisible characters, fake system messages, "
                           "instructions to change the verdict.", size=14, space=12)
                bullet(tf, "Nothing was found.", size=16, bold=True, color=t["good"])
            else:
                for key, title in (("deck", "Pitch deck"), ("web_sources", "Web sources")):
                    findings = (sec.get(key) or {}).get("findings") or []
                    if not findings:
                        continue
                    bullet(tf, title, size=13, bold=True, color=t["muted"], space=4)
                    for f in findings[:6]:
                        color = {"critical": t["bad"], "high": t["bad"],
                                 "medium": t["warn"]}.get(f.get("severity"), t["muted"])
                        bullet(tf, f"[{f.get('severity')}] {f.get('where')} — "
                                   f"{str(f.get('detail'))[:150]}",
                               size=11.5, color=color, space=3, level=1)
                    bullet(tf, "", size=6, space=8)

        # -------------------------------------------------------- endnote
        s = new_slide("Method and caveats")
        tf = body_box(s)
        bullet(tf, "How this was produced", size=13, bold=True, color=t["muted"], space=6)
        for line in (
            f"Deck extraction, independent market research, and comparison were run as "
            f"three separate passes so the market view is not anchored on the deck's claims.",
            f"Model: {h['model']}. Research: {h['research']}.",
            f"Analytical lens: {h['lens']}.",
            f"Generated {h['generated']}.",
        ):
            bullet(tf, f"·  {line}", size=12.5, space=5)
        bullet(tf, "", size=6, space=10)
        bullet(tf, "AI-generated analysis. Verify every figure against the cited source "
                   "before relying on it. Not investment advice.",
               size=12, bold=True, color=t["warn"])

        path = out_dir / f"{base}_{lens}.pptx"
        prs.save(str(path))
        paths.append(str(path))
    return paths
