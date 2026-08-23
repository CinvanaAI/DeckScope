"""Turn a pitch deck in any common format into slide-numbered plain text.

Supported: .pptx .pdf .docx .txt .md .html .json, plus raw text and http(s) URLs.
Scanned/image-only PDFs are detected and reported rather than silently returning
an empty deck.
"""
from __future__ import annotations

import json
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

SUPPORTED_EXTENSIONS = {".pptx", ".ppt", ".pdf", ".docx", ".txt", ".md",
                        ".markdown", ".html", ".htm", ".json"}


class DeckLoadError(RuntimeError):
    """The deck could not be read into usable text."""


@dataclass
class DeckDocument:
    text: str
    n_slides: int
    source: str
    fmt: str
    warnings: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    #: An on-disk copy of the ORIGINAL file, when the deck came from a URL.
    #:
    #: File-level forensics — hidden slides, speaker notes, off-slide shapes,
    #: invisible PDF text, document metadata — can only run against the real
    #: binary. Remote decks used to have their temporary download deleted the
    #: instant the text was extracted, and the URL was handed to the scanner
    #: instead. A URL is not a file, so the scanner found nothing and reported
    #: nothing: every forensic protection the docs advertise silently did not
    #: apply to any deck fetched over the network, which is exactly the deck most
    #: likely to be hostile.
    #:
    #: The caller owns this file and must call `cleanup()` when done.
    local_path: Optional[str] = None

    def cleanup(self) -> None:
        """Delete the temporary original, if we made one."""
        if not self.local_path:
            return
        try:
            Path(self.local_path).unlink()
        except OSError:
            pass
        finally:
            self.local_path = None

    @property
    def is_thin(self) -> bool:
        """Too little text to analyze — usually an image-only deck."""
        return len(self.text.strip()) < 400


def load_deck(path_or_text: str, *, is_text: bool = False) -> DeckDocument:
    if is_text:
        return _from_text(path_or_text, source="inline text")
    if re.match(r"^https?://", path_or_text.strip(), re.I):
        return _from_url(path_or_text.strip())

    p = Path(path_or_text).expanduser()
    if not p.exists():
        raise DeckLoadError(
            f"No file at {p}. Check the path, or drag the file onto the app window."
        )
    ext = p.suffix.lower()
    loaders = {
        ".pptx": _from_pptx, ".ppt": _from_pptx, ".pdf": _from_pdf,
        ".docx": _from_docx, ".json": _from_json,
        ".html": _from_html, ".htm": _from_html,
    }
    if ext in loaders:
        return loaders[ext](p)
    if ext in {".txt", ".md", ".markdown"} or ext == "":
        return _from_text(p.read_text(encoding="utf-8", errors="replace"), source=str(p),
                          fmt=ext.lstrip(".") or "txt")
    raise DeckLoadError(
        f"Cannot read {ext or 'that file'}. Supported: "
        f"{', '.join(sorted(SUPPORTED_EXTENSIONS))}. "
        f"Tip: export the deck to PDF and try again."
    )


# ---------------------------------------------------------------- formats

#: Caps on what an Office file may expand to. A .pptx and a .docx are both zip
#: archives, and a zip archive can be a few kilobytes on disk that expands to
#: gigabytes in memory. python-pptx and python-docx expand what they are given
#: without asking, so the check has to happen before they are handed the file —
#: at which point a malicious deck is a memory-exhaustion primitive rather than
#: a document. The limits are far above any real deck: a heavily illustrated
#: 200-slide deck is tens of megabytes uncompressed.
MAX_OFFICE_UNCOMPRESSED = 512 * 1024 * 1024   # 512 MB across all members
MAX_OFFICE_RATIO = 200                        # expanded : on-disk
MAX_OFFICE_MEMBERS = 10_000


def _check_office_archive(p: Path) -> None:
    """Refuse an Office file that expands out of proportion to its size.

    Reads only the zip central directory, so this costs nothing on a real deck
    and never expands the payload it is judging.
    """
    import zipfile

    try:
        with zipfile.ZipFile(p) as zf:
            infos = zf.infolist()
            if len(infos) > MAX_OFFICE_MEMBERS:
                raise DeckLoadError(
                    f"This file contains {len(infos):,} internal parts, far more than "
                    f"a real presentation or document. Refusing to open it.")
            total = sum(i.file_size for i in infos)
            if total > MAX_OFFICE_UNCOMPRESSED:
                raise DeckLoadError(
                    f"This file expands to {total / 1e6:,.0f} MB, past the "
                    f"{MAX_OFFICE_UNCOMPRESSED / 1e6:,.0f} MB limit. Refusing to open "
                    f"it. If the deck is genuinely this large, export it to PDF.")
            on_disk = p.stat().st_size or 1
            if total / on_disk > MAX_OFFICE_RATIO and total > 32 * 1024 * 1024:
                raise DeckLoadError(
                    f"This file is {on_disk / 1e6:.1f} MB on disk but expands to "
                    f"{total / 1e6:,.0f} MB — a {total / on_disk:,.0f}x ratio that no "
                    f"ordinary deck produces. Refusing to open it.")
    except zipfile.BadZipFile:
        raise DeckLoadError(
            f"{p.name} is not a readable Office file — the archive is damaged or it "
            f"is not really a .pptx/.docx. Try re-exporting it.") from None


SLIDE_MARKER = re.compile(r"^\s*-{2,}\s*(?:slide|page)\s*(\d+)\s*-{2,}\s*$",
                          re.I | re.M)


def _from_text(text: str, source: str, fmt: str = "text") -> DeckDocument:
    marked = SLIDE_MARKER.findall(text)
    if len(marked) > 1:
        return DeckDocument(text, len(marked), source, fmt)
    chunks = re.split(r"\n(?:---+|===+|\f)\n", text)
    if len(chunks) > 1:
        body = "\n\n".join(f"--- Slide {i} ---\n{c.strip()}"
                           for i, c in enumerate(chunks, 1) if c.strip())
        return DeckDocument(body, len(chunks), source, fmt)
    return DeckDocument(text, max(1, text.count("\f") + 1), source, fmt)


def _from_pptx(p: Path) -> DeckDocument:
    try:
        from pptx import Presentation
    except ImportError:
        raise DeckLoadError("Reading .pptx needs python-pptx: pip install python-pptx") from None
    if p.suffix.lower() == ".ppt":
        raise DeckLoadError("Legacy .ppt isn't supported. Save as .pptx or export to PDF.")
    # Before python-pptx expands anything.
    _check_office_archive(p)
    prs = Presentation(str(p))
    out, warnings, images = [], [], 0
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                images += 1
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" | ".join(c.text.strip() for c in row.cells))
            if getattr(shape, "has_chart", False):
                try:
                    cats = [str(c) for c in shape.chart.plots[0].categories]
                    lines.append(f"[chart categories: {', '.join(cats)}]")
                except Exception:  # noqa: BLE001
                    lines.append("[chart present, values not extractable]")
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                lines.append(f"[speaker notes] {slide.notes_slide.notes_text_frame.text.strip()}")
        except Exception:  # noqa: BLE001
            pass
        out.append("\n".join(lines))
    if images and sum(len(o) for o in out) < 1500:
        warnings.append(
            f"{images} images found but very little text — key numbers may live inside "
            f"graphics this reader cannot see."
        )
    return DeckDocument("\n\n".join(out), len(prs.slides), str(p), "pptx", warnings)


def _from_pdf(p: Path) -> DeckDocument:
    text_pages: List[str] = []
    warnings: List[str] = []
    try:
        import pdfplumber

        with pdfplumber.open(str(p)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                t = page.extract_text() or ""
                tables = []
                try:
                    for tbl in page.extract_tables() or []:
                        tables.append("\n".join(
                            " | ".join((c or "").strip() for c in row) for row in tbl))
                except Exception:  # noqa: BLE001
                    pass
                block = f"--- Slide {i} ---\n{t}"
                if tables:
                    block += "\n[tables]\n" + "\n\n".join(tables)
                text_pages.append(block)
    except ImportError:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise DeckLoadError("Reading PDFs needs pdfplumber or pypdf: "
                                "pip install pdfplumber") from None
        reader = PdfReader(str(p))
        for i, page in enumerate(reader.pages, 1):
            text_pages.append(f"--- Slide {i} ---\n{page.extract_text() or ''}")
    except Exception as exc:  # noqa: BLE001
        raise DeckLoadError(f"Could not read {p.name}: {exc}") from None

    body = "\n\n".join(text_pages)
    if len(re.sub(r"--- Slide \d+ ---", "", body).strip()) < 300:
        warnings.append(
            "This PDF appears to be scanned or image-only — almost no selectable text. "
            "Run it through OCR first, or export the original deck to PDF rather than "
            "printing it to images."
        )
    return DeckDocument(body, len(text_pages), str(p), "pdf", warnings)


def _from_docx(p: Path) -> DeckDocument:
    try:
        import docx
    except ImportError:
        raise DeckLoadError("Reading .docx needs python-docx: pip install python-docx") from None
    # Before python-docx expands anything.
    _check_office_archive(p)
    d = docx.Document(str(p))
    parts = [para.text for para in d.paragraphs if para.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return DeckDocument("\n".join(parts), 1, str(p), "docx")


def _from_html(p: Path) -> DeckDocument:
    raw = p.read_text(encoding="utf-8", errors="replace")
    return _from_text(_strip_html(raw), str(p), "html")


def _from_json(p: Path) -> DeckDocument:
    data = json.loads(p.read_text(encoding="utf-8"))
    if isinstance(data, list):
        body = "\n\n".join(f"--- Slide {i} ---\n{json.dumps(s, indent=2)}"
                           for i, s in enumerate(data, 1))
        return DeckDocument(body, len(data), str(p), "json")
    return DeckDocument(json.dumps(data, indent=2), 1, str(p), "json")


def _from_url(url: str) -> DeckDocument:
    """Fetch and read a deck from a URL, with SSRF and size guards.

    Every safety decision lives in ingest/fetch.py; this function only turns the
    fetched bytes into a document. Binary formats are written to a uniquely named
    temporary file — the earlier fixed filename meant two concurrent panel runs
    could overwrite each other's download mid-read.
    """
    import tempfile

    from .fetch import FetchError, fetch_url

    try:
        got = fetch_url(url)
    except FetchError as exc:
        raise DeckLoadError(str(exc)) from None

    if got.suffix in (".pdf", ".pptx", ".docx"):
        fd, tmp_name = tempfile.mkstemp(prefix="deckscope_", suffix=got.suffix)
        tmp = Path(tmp_name)
        try:
            with os.fdopen(fd, "wb") as fh:
                fh.write(got.content)
            reader = {".pdf": _from_pdf, ".pptx": _from_pptx, ".docx": _from_docx}
            doc = reader[got.suffix](tmp)
        except Exception:
            try:
                tmp.unlink()
            except OSError:
                pass
            raise
        # Deliberately NOT deleted here. Forensics needs the original bytes, and
        # deleting them at this point is what made every file-level check skip
        # remote decks. Ownership passes to the caller via `local_path`.
        doc.local_path = str(tmp)
        doc.source = got.final_url
        doc.metadata["fetched_from"] = url
        return doc

    text = got.content.decode("utf-8", "replace")
    if got.suffix == ".html":
        text = _strip_html(text)
    doc = _from_text(text, got.final_url, got.suffix.lstrip(".") or "text")
    doc.metadata["fetched_from"] = url
    return doc


def _strip_html(raw: str) -> str:
    raw = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.S | re.I)
    raw = re.sub(r"<br\s*/?>|</p>|</div>|</li>", "\n", raw, flags=re.I)
    raw = re.sub(r"<[^>]+>", " ", raw)
    import html as _html
    return re.sub(r"\n{3,}", "\n\n", _html.unescape(raw))
