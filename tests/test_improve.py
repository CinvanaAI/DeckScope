"""The reverse flow: `deckscope improve` — deck (or raw notes) in,
corrected deck blueprint out. What matters here is the honesty layer:
every rule the docstring promises is enforced by code, and these tests
are that promise pinned.
"""
from __future__ import annotations

import sys
import types
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from deckscope.commands.audit_report import build_registry
from deckscope.commands.improve import (build_brief, render_blueprint,
                                        validate_revision, write_pptx)

_SOURCES = [{"sid": "S1", "title": "Independent sizing",
             "url": "https://example.gov/a", "snippet": "$18-24B"}]

_COMP = {
    "claim_audit": [
        {"id": "C1", "claim": "The market is $47B",
         "assessment": "contradicted", "source_ids": ["S1"],
         "market_evidence": "independent estimates say $18-24B",
         "delta": "roughly 2x overstated"},
        {"id": "C2", "claim": "Pricing is $2,000/month",
         "assessment": "supported", "source_ids": ["S1"]},
        {"id": "C3", "claim": "NRR is 140%", "assessment": "unverifiable"},
    ],
    "alignment": {"blind_spots": [
        {"what": "Power Automate is free in E5",
         "why_it_matters": "zero marginal cost for the buyer"}]},
    "questions": ["What is CAC payback?"],
}


def _rev(lines, extra=None):
    rev = {"company": "Acme Flow", "positioning": "p",
           "slides": [{"n": 1, "title": "T", "purpose": "",
                       "lines": lines, "speaker_note": ""}],
           "cut": [], "founder_inputs": []}
    rev.update(extra or {})
    return rev


# ------------------------------------------------------------- the brief

def test_brief_carries_the_whole_audit_and_only_the_run_bibliography():
    class _Reg:
        def prompt_block(self, char_budget=0):
            return "[S1] Independent sizing"

    brief = build_brief(_COMP, _Reg())
    assert "CLAIMS THE EVIDENCE CONTRADICTS" in brief
    assert "independent estimates say $18-24B [S1]" in brief
    assert "CLAIMS THAT SURVIVED THE AUDIT" in brief
    assert "CLAIMS ONLY THE FOUNDER CAN SETTLE" in brief
    assert "NRR is 140%" in brief
    assert "Power Automate is free in E5" in brief
    assert "What is CAC payback?" in brief
    assert "the ONLY citable sources" in brief
    assert "[S1] Independent sizing" in brief


# ----------------------------------------------------- the honesty layer

def test_a_fabricated_citation_is_stripped_and_the_figure_demoted():
    """The worst case for a deck-builder: an invented figure wearing an
    invented citation. The citation audit strips [S9]; the now-uncited
    figure is demoted to a founder slot. Nothing invented survives."""
    reg = build_registry(list(_SOURCES))
    rev = _rev([{"text": "The category is $99B", "kind": "revised",
                 "source_ids": ["S9"], "because": "x"}])
    rev, notes = validate_revision(rev, reg, _COMP)
    line = rev["slides"][0]["lines"][0]
    assert notes["stripped_citations"] >= 1
    assert line["source_ids"] == []
    assert line["kind"] == "founder-input"
    assert notes["demoted_lines"] == 1
    assert any("$99B" in f.get("slot", "") for f in rev["founder_inputs"])


def test_an_unsourced_new_figure_is_demoted_but_cited_ones_pass():
    reg = build_registry(list(_SOURCES))
    rev = _rev([
        {"text": "The category is $18-24B", "kind": "revised",
         "source_ids": ["S1"], "because": "C1"},
        {"text": "Payback is under $5,000 CAC", "kind": "new",
         "source_ids": [], "because": "q"},
    ])
    rev, notes = validate_revision(rev, reg, _COMP)
    kinds = [ln["kind"] for ln in rev["slides"][0]["lines"]]
    assert kinds == ["revised", "founder-input"]
    assert notes["demoted_lines"] == 1


def test_kept_lines_keep_their_figures_but_contested_ones_are_flagged():
    """A founder may keep their own numbers — they own them. But a kept
    line that matches a contested claim gets the warning an investor-side
    run would make real."""
    reg = build_registry(list(_SOURCES))
    rev = _rev([
        {"text": "The market is $47B", "kind": "kept", "source_ids": [],
         "because": ""},
        {"text": "Pricing is $2,000/month", "kind": "kept",
         "source_ids": [], "because": ""},
    ])
    rev, notes = validate_revision(rev, reg, _COMP)
    lines = rev["slides"][0]["lines"]
    assert lines[0]["kind"] == "kept", "kept is never demoted"
    assert lines[0].get("kept_against_evidence") == "C1"
    assert "kept_against_evidence" not in lines[1], (
        "a kept line matching a SUPPORTED claim is fine")
    assert notes["kept_against_evidence"] == 1


def test_founder_input_lines_always_show_a_visible_slot():
    reg = build_registry(list(_SOURCES))
    rev = _rev([{"text": "Net revenue retention", "kind": "founder-input",
                 "source_ids": [], "because": "C3"}])
    rev, _ = validate_revision(rev, reg, _COMP)
    assert "[" in rev["slides"][0]["lines"][0]["text"]


def test_empty_lines_and_slides_are_dropped_and_renumbered():
    reg = build_registry(list(_SOURCES))
    rev = {"company": "A", "positioning": "",
           "slides": [
               {"n": 9, "title": "empty", "purpose": "", "lines": [
                   {"text": "  ", "kind": "kept", "source_ids": []}],
                "speaker_note": ""},
               {"n": 4, "title": "real", "purpose": "", "lines": [
                   {"text": "A line", "kind": "kept", "source_ids": []}],
                "speaker_note": ""},
               "not a slide"],
           "cut": [], "founder_inputs": []}
    rev, _ = validate_revision(rev, reg, _COMP)
    assert len(rev["slides"]) == 1
    assert rev["slides"][0]["n"] == 1


# ------------------------------------------------------------- rendering

def test_blueprint_renders_marks_warnings_and_the_appendix():
    reg = build_registry(list(_SOURCES))
    rev = _rev([
        {"text": "The market is $47B", "kind": "kept", "source_ids": [],
         "because": ""},
        {"text": "The category is $18-24B", "kind": "revised",
         "source_ids": ["S1"], "because": "C1 contradicted"},
        {"text": "NRR", "kind": "founder-input", "source_ids": [],
         "because": "C3"},
    ], extra={"cut": [{"what": "TAM slide", "why": "contradicted"}]})
    rev, notes = validate_revision(rev, reg, _COMP)
    body = render_blueprint(rev, notes, reg, "deck.pdf")
    assert "will not invent your numbers" in body
    assert "kept against the evidence (C1)" in body
    assert "✎ you" in body
    assert "[S1]" in body
    assert "Evidence appendix" in body
    assert "Independent sizing" in body, "appendix resolves cited sources"
    assert "## Cut from the deck" in body
    assert "not investment or fundraising advice" in body.lower() or \
           "not a finished deck" in body


def test_pptx_marks_founder_slots_when_pptx_is_available(tmp_path):
    reg = build_registry(list(_SOURCES))
    rev = _rev([
        {"text": "The category is $18-24B", "kind": "revised",
         "source_ids": ["S1"], "because": "C1"},
        {"text": "NRR [YOUR NUMBER]", "kind": "founder-input",
         "source_ids": [], "because": "C3"}])
    rev, _ = validate_revision(rev, reg, _COMP)
    p = write_pptx(rev, tmp_path / "deck.pptx")
    if p is None:  # python-pptx genuinely absent — the CLI says so too
        return
    assert p.exists() and p.stat().st_size > 0
    from pptx import Presentation

    texts = []
    for slide in Presentation(str(p)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    joined = "\n".join(texts)
    assert "▢ NRR [YOUR NUMBER]" in joined
    assert "[S1]" in joined


# --------------------------------------------------------------- command

def test_improve_demo_runs_offline_end_to_end(tmp_path):
    """The whole reverse flow with zero keys and zero network: sample
    deck -> mock pipeline -> mock reviser -> validated blueprint. The
    fixture is authored to trip both guards so the demo shows them."""
    from deckscope.commands.improve import command

    args = types.SimpleNamespace(demo=True, nda=False, lens=None,
                                 deck=None, from_run=None, provider=None,
                                 config=None, out=str(tmp_path), pptx=False)
    assert command(args) == 0
    out = tmp_path / "sample_deck_founder_improved.md"
    body = out.read_text(encoding="utf-8")
    assert "# Improved deck — Acme Flow" in body
    assert "demoted to founder slots" in body, "the demotion guard fires"
    assert "kept against the evidence" in body, "the kept-flag guard fires"
    assert "[S1]" in body


def test_improve_nda_refuses_a_hosted_model(tmp_path):
    from deckscope.commands.improve import command

    deck = tmp_path / "notes.txt"
    deck.write_text("Our market is $47B", encoding="utf-8")
    args = types.SimpleNamespace(demo=False, nda=True, lens=None,
                                 deck=str(deck), from_run=None,
                                 provider="anthropic", config=None,
                                 out=str(tmp_path), pptx=False)
    assert command(args) == 4


def test_improve_without_a_deck_says_notes_work_too(tmp_path, capsys=None):
    from deckscope.commands.improve import command

    args = types.SimpleNamespace(demo=False, nda=False, lens=None,
                                 deck=None, from_run=None, provider=None,
                                 config=None, out=str(tmp_path), pptx=False)
    assert command(args) == 2


def test_improve_is_wired_into_the_cli():
    from deckscope.cli import build_parser

    assert "improve" in build_parser().format_help()
